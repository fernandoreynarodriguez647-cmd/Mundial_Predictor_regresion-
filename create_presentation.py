from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import json
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Color Palette ───
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x0A, 0x1F, 0x44)
NAVY_LIGHT = RGBColor(0x12, 0x2D, 0x5E)
SKY_BLUE = RGBColor(0x5D, 0x9E, 0xCF)
SKY_BLUE_LIGHT = RGBColor(0xA8, 0xCC, 0xE8)
SOFT_WHITE = RGBColor(0xF2, 0xF5, 0xFA)
DARK_TEXT = RGBColor(0x1A, 0x2A, 0x44)
MEDIUM_TEXT = RGBColor(0x3D, 0x5A, 0x7A)
LIGHT_TEXT = RGBColor(0x7A, 0x95, 0xB0)
ACCENT_GOLD = RGBColor(0xC8, 0xA2, 0x4A)
ACCENT_TEAL = RGBColor(0x2E, 0x86, 0xAB)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE_BG = RGBColor(0xE8, 0xEE, 0xF6)
SECTION_BG = RGBColor(0xF7, 0xF9, 0xFC)

# ─── Helper functions ───
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, color, border=None, border_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color, border=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border)
    else:
        shape.line.fill.background()
    # adjust corner radius
    sp = shape._element
    spPr = sp.find(qn('a:spPr'))
    if spPr is not None:
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = prstGeom.makeelement(qn('a:avLst'), {})
                prstGeom.append(avLst)
            gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val 15000'})
            avLst.append(gd)
    return shape

def tb(slide, left, top, width, height, text, size=16, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, name='Calibri', anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = align
    except:
        pass
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    return box

def add_para(tf, text, size=14, color=MEDIUM_TEXT, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    p.space_before = space_before
    return p

def slide_header(slide, title, subtitle=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.15), NAVY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), prs.slide_width, Inches(0.04), SKY_BLUE)
    tb(slide, Inches(0.8), Inches(0.18), Inches(11), Inches(0.6), title, size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        tb(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.35), subtitle, size=13, color=SKY_BLUE_LIGHT, anchor=MSO_ANCHOR.MIDDLE)

def card(slide, left, top, width, height, title, body_lines, title_color=NAVY, body_color=MEDIUM_TEXT):
    add_rounded_rect(slide, left, top, width, height, CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04), SKY_BLUE)
    tb(slide, left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.35), title, size=16, color=title_color, bold=True)
    y_offset = top + Inches(0.6)
    for line in body_lines:
        tb(slide, left + Inches(0.25), y_offset, width - Inches(0.5), Inches(0.28), line, size=12, color=body_color)
        y_offset += Inches(0.26)

def icon_bullet(slide, left, top, width, icon_text, label, desc, label_color=NAVY, desc_color=MEDIUM_TEXT):
    tb(slide, left, top, Inches(0.3), Inches(0.25), icon_text, size=12, color=SKY_BLUE, bold=True)
    tb(slide, left + Inches(0.35), top, Inches(2), Inches(0.25), label, size=13, color=label_color, bold=True)
    tb(slide, left + Inches(0.35), top + Inches(0.25), Inches(width - 0.35), Inches(0.25), desc, size=11, color=desc_color)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – PORTADA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
# Decorative top bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06), SKY_BLUE)
# Bottom decoration
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), prs.slide_width, Inches(0.06), SKY_BLUE)
# Left accent bar
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.5), Inches(0.06), Inches(3.5), SKY_BLUE)
# Title
tb(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(0.8), "Predictor Mundial FIFA 2026", size=42, color=WHITE, bold=True)
tb(slide, Inches(1.2), Inches(3.4), Inches(10), Inches(0.5), "Sistema de prediccion deportiva basado en Machine Learning", size=18, color=SKY_BLUE_LIGHT)
tb(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.4), "Pipeline progresivo · Modelos de clasificacion y regresion · Dashboard interactivo", size=13, color=LIGHT_TEXT)
# Bottom info
tb(slide, Inches(1.2), Inches(5.5), Inches(5), Inches(0.3), "Ciencia de Datos  |  Inteligencia Artificial  |  Deportes", size=12, color=SKY_BLUE_LIGHT)
tb(slide, Inches(1.2), Inches(5.9), Inches(5), Inches(0.3), "Fases eliminatorias  |  31 features  |  Ensemble ponderado", size=12, color=LIGHT_TEXT)
# Right decorative circle
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.5), Inches(2.0), Inches(2.5), Inches(2.5), RGBColor(0x0F, 0x2B, 0x55))
add_shape(slide, MSO_SHAPE.OVAL, Inches(10.8), Inches(2.3), Inches(1.9), Inches(1.9), RGBColor(0x14, 0x35, 0x68))

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Agenda")
items = [
    ("01", "Introduccion", "Motivacion, objetivo y alcance del proyecto"),
    ("02", "Datos", "Fuentes de informacion y estructura de datos"),
    ("03", "Ingenieria de Features", "31 variables predictivas disenadas"),
    ("04", "Modelos", "Clasificadores, regresores y ensemble ponderado"),
    ("05", "Pipeline", "Flujo progresivo R32 a FINAL con validacion"),
    ("06", "Resultados Globales", "Accuracy y resumen por fase"),
    ("07", "Resultados por Fase", "R32, R16, QF y SF detallados"),
    ("08", "Feature Importance", "Variables mas influyentes"),
    ("09", "Dashboard", "Visualizacion interactiva con Streamlit"),
    ("10", "Evaluacion", "Metricas y validacion automatica"),
    ("11", "Conclusiones", "Hallazgos, logros y trabajo futuro"),
]
for i, (num, title, desc) in enumerate(items):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.5 + row * 2.8)
    add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.4), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(2.9), Inches(0.04), SKY_BLUE)
    tb(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.4), num, size=22, color=SKY_BLUE, bold=True)
    tb(slide, x + Inches(0.7), y + Inches(0.25), Inches(2), Inches(0.4), title, size=15, color=NAVY, bold=True)
    tb(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.5), Inches(1.2), desc, size=12, color=MEDIUM_TEXT)
# 3rd row for items 9-11
if len(items) > 8:
    for i in range(8, len(items)):
        col = (i - 8) % 4
        row = 2
        x = Inches(0.6 + col * 3.15)
        y = Inches(1.5 + row * 2.8)
        num, title, desc = items[i]
        add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.4), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(2.9), Inches(0.04), SKY_BLUE)
        tb(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.4), num, size=22, color=SKY_BLUE, bold=True)
        tb(slide, x + Inches(0.7), y + Inches(0.25), Inches(2), Inches(0.4), title, size=15, color=NAVY, bold=True)
        tb(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.5), Inches(1.2), desc, size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – INTRODUCCION
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Introduccion", "Motivacion, objetivo y alcance del proyecto")
# Card 1
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.3), "Objetivo", size=17, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(2.1), Inches(5.3), Inches(1.8),
   "Desarrollar un sistema reproducible de prediccion\ndeportiva basado en Machine Learning para las\nfases eliminatorias de la Copa Mundial FIFA 2026.\n\nEl pipeline avanza fase por fase utilizando\nunicamente resultados oficiales ya disputados,\ngarantizando consistencia en cada etapa.", size=12, color=MEDIUM_TEXT)
# Card 2
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.3), "Motivacion", size=17, color=NAVY, bold=True)
tb(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(1.8),
   "Demostrar aplicacion de Machine Learning en\nanalisis deportivo con un enfoque progresivo.\n\nCombinacion de modelos de clasificacion para\npredecir el ganador, regresion para el marcador\ny modelo especializado para definicion por penales.", size=12, color=MEDIUM_TEXT)
# Card 3
add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.3), "Arquitectura", size=17, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(5.1), Inches(5.3), Inches(1.6),
   "Datos historicos (2014, 2018, 2022)  >  Feature Engineering\n>  Entrenamiento de modelos  >  Prediccion por fase\n(R32  >  R16  >  QF  >  SF  >  FINAL)\n\nCada fase congela los enfrentamientos oficiales.", size=12, color=MEDIUM_TEXT)
# Card 4
add_rounded_rect(slide, Inches(6.9), Inches(4.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(4.7), Inches(5.5), Inches(0.3), "Alcance", size=17, color=NAVY, bold=True)
tb(slide, Inches(7.2), Inches(5.1), Inches(5.3), Inches(1.6),
   "6 fases eliminatorias del Mundial 2026\n3 modelos de clasificacion + 2 regresores + penal\nDashboard interactivo en Streamlit\nReportes automaticos y metricas por fase\nValidacion contra resultados reales", size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – STACK TECNOLOGICO
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Stack Tecnologico")
techs = [
    ("Python", "Lenguaje principal", "Pandas, NumPy, Scikit-learn"),
    ("XGBoost", "Gradient Boosting", "Clasificacion y regresion"),
    ("Streamlit", "Dashboard interactivo", "Visualizacion de datos"),
    ("Plotly", "Graficos interactivos", "Curvas y feature importance"),
    ("Joblib", "Persistencia", "Modelos guardados en disco"),
    ("Pytest", "Testing", "Validacion del pipeline"),
]
for i, (name, role, detail) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.5 + row * 2.9)
    add_rounded_rect(slide, x, y, Inches(3.9), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(3.9), Inches(0.04), SKY_BLUE)
    tb(slide, x + Inches(0.25), y + Inches(0.2), Inches(3.4), Inches(0.35), name, size=18, color=NAVY, bold=True)
    tb(slide, x + Inches(0.25), y + Inches(0.6), Inches(3.4), Inches(0.25), role, size=13, color=SKY_BLUE)
    tb(slide, x + Inches(0.25), y + Inches(1.0), Inches(3.4), Inches(1.2), detail, size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – DATOS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Datos", "Fuentes de informacion y estructura")
# Structure diagram
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(1.0), NAVY)
tb(slide, Inches(0.7), Inches(1.55), Inches(3.6), Inches(0.3), "data/", size=16, color=WHITE, bold=True)
tb(slide, Inches(0.7), Inches(1.9), Inches(3.6), Inches(0.5), "raw/   >   Datos originales sin procesar\nprocessed/   >   Datos transformados por fase", size=11, color=SKY_BLUE_LIGHT)
# Files as cards
files = [
    ("historical_matches.csv", "Partidos eliminatorios 2014-2022", "36 registros"),
    ("teams_static.csv", "Estadisticas de selecciones", "Ranking, grupos"),
    ("players.csv", "Jugadores clave", "Edad, goles, lesiones"),
    ("historical_h2h.csv", "Head-to-Head historico", "Enfrentamientos"),
    ("bracket_map.json", "Cuadro eliminatorio", "Estructura del torneo"),
    ("official_results.csv", "Resultados oficiales", "Actualizados"),
]
for i, (name, desc, detail) in enumerate(files):
    col = i % 3
    row = i // 3
    x = Inches(4.9 + col * 2.85)
    y = Inches(1.5 + row * 1.9)
    add_rounded_rect(slide, x, y, Inches(2.65), Inches(1.65), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(2.65), Inches(0.03), SKY_BLUE)
    tb(slide, x + Inches(0.15), y + Inches(0.12), Inches(2.35), Inches(0.3), name, size=11, color=SKY_BLUE, bold=True)
    tb(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.35), Inches(0.25), desc, size=12, color=NAVY, bold=True)
    tb(slide, x + Inches(0.15), y + Inches(0.85), Inches(2.35), Inches(0.6), detail, size=11, color=MEDIUM_TEXT)
# Bottom note
add_rounded_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.03), SKY_BLUE)
tb(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.3), "Procesamiento", size=15, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.7),
   "Los datos crudos se transforman mediante ingenieria de features para generar 31 variables predictoras.\nCada fase del torneo produce sus propios archivos de prediccion, comparacion y metricas en data/processed/.", size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – FEATURE ENGINEERING
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Ingenieria de Features", "31 variables predictivas organizadas en 6 categorias")
cats = [
    ("Ranking y Mundial", "fifa_ranking_pts_diff\nfifa_ranking_pos_diff\nwc_appearances_diff\nwc_best_stage_diff\nwc_titles_diff"),
    ("Rendimiento Grupos", "goals_for_group_diff\ngoals_against_group_diff\ngroup_strength_score_diff\ngroup_position_a / b"),
    ("Jugadores Clave", "key_player_availability_diff\ntop_scorer_goals_tournament_diff\nkey_players_avg_age_diff\nkey_players_prev_wc_goals_diff"),
    ("Estadisticas", "possession_avg_diff\nshots_on_target_avg_diff\ndefensive_efficiency_diff\nclean_sheets_ko_diff"),
    ("Head-to-Head", "h2h_wins_a / h2h_wins_b\nh2h_draws\nh2h_goal_diff_avg\nh2h_penalties_a_win_rate"),
    ("Historial KO", "ko_stage_historical_win_rate_a\nko_stage_historical_win_rate_b\nstage_ordinal\nhost_advantage"),
]
for i, (cat, feats) in enumerate(cats):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.5 + row * 2.9)
    add_rounded_rect(slide, x, y, Inches(3.9), Inches(2.6), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(3.9), Inches(0.04), SKY_BLUE)
    tb(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.3), cat, size=15, color=NAVY, bold=True)
    tb(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.5), Inches(1.8), feats, size=11, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 – MODELOS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Modelos de Machine Learning", "Clasificacion, regresion y ensemble ponderado")
# Classifiers
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.3), "Clasificadores", size=17, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(2.1), Inches(5.3), Inches(1.8),
   "Logistic Regression  |  max_iter=2000, class_weight=balanced\nRandom Forest       |  500 arboles, max_depth=5\nXGBoost              |  400 arboles, learning_rate=0.03\n\nCada modelo predice P(gana Team_A) para cada partido.\nOutput: probabilidad de avance por equipo.", size=12, color=MEDIUM_TEXT)
# Regressors
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.3), "Regresores y Penal", size=17, color=NAVY, bold=True)
tb(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(1.8),
   "XGBRegressor (score_a)  |  objective=count:poisson\nXGBRegressor (score_b)  |  objective=count:poisson\n\nPredicen goles en tiempo regular (90 min) por equipo.\nSi hay empate: LogisticRegression para penales\n(solo si hay 6+ partidos con penales en training).", size=12, color=MEDIUM_TEXT)
# Ensemble
add_rounded_rect(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(0.3), "Ensemble Ponderado", size=17, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(5.1), Inches(5.3), Inches(1.6),
   "Cada clasificador recibe un peso dinamico basado en\nsu accuracy historica en fases anteriores.\n\nPesos iniciales: 1/3 cada uno.\nSe recalculan tras cada validacion (update_stage)\nusando la formula: Peso_m = acc_m / sum(acc_total)", size=12, color=MEDIUM_TEXT)
# Training
add_rounded_rect(slide, Inches(6.9), Inches(4.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(4.7), Inches(5.5), Inches(0.3), "Entrenamiento", size=17, color=NAVY, bold=True)
tb(slide, Inches(7.2), Inches(5.1), Inches(5.3), Inches(1.6),
   "Los modelos se entrenan con datos de fases KO de\nMundiales 2014, 2018 y 2022 (36 partidos).\n\nSe guardan en saved_models/ con Joblib para\nreutilizarse entre ejecuciones.\nReentrenables con: python scripts/train_models.py", size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 – PIPELINE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Pipeline Progresivo", "Flujo de ejecucion fase por fase")
# Stage boxes with arrows
stages = [
    ("R32", "Dieciseisavos"),
    ("R16", "Octavos"),
    ("QF", "Cuartos"),
    ("SF", "Semifinales"),
    ("F", "Final"),
]
x = 0.4
for stage, name in stages:
    add_rounded_rect(slide, Inches(x), Inches(1.5), Inches(2.3), Inches(1.5), NAVY)
    tb(slide, Inches(x + 0.1), Inches(1.55), Inches(2.1), Inches(0.4), stage, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(slide, Inches(x + 0.1), Inches(2.05), Inches(2.1), Inches(0.3), name, size=13, color=SKY_BLUE_LIGHT, align=PP_ALIGN.CENTER)
    if stage != "F":
        tb(slide, Inches(x + 2.3), Inches(1.9), Inches(0.3), Inches(0.3), ">", size=20, color=SKY_BLUE, bold=True, align=PP_ALIGN.CENTER)
    x += 2.6
# Flow
add_rounded_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.03), SKY_BLUE)
tb(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.3), "Flujo de ejecucion", size=16, color=NAVY, bold=True)
# Steps
steps = [
    ("python scripts/run_stage.py <FASE>", "Predice los partidos de la fase actual usando modelos entrenados con datos historicos"),
    ("python scripts/update_stage.py <FASE>", "Valida las predicciones contra resultados oficiales y recalcula pesos del ensemble"),
    ("Repetir", "Avanzar a la siguiente fase hasta llegar a la Final"),
]
y = Inches(3.95)
for cmd, desc in steps:
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.04), Inches(0.6), SKY_BLUE)
    tb(slide, Inches(1.1), y, Inches(4), Inches(0.6), cmd, size=11, color=ACCENT_TEAL, bold=True)
    tb(slide, Inches(5.2), y, Inches(7.2), Inches(0.6), desc, size=12, color=MEDIUM_TEXT)
    y += Inches(0.7)
# Guard
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.04), Inches(0.6), ACCENT_GOLD)
tb(slide, Inches(1.1), y, Inches(4), Inches(0.6), "Mecanismo de guardia", size=11, color=ACCENT_GOLD, bold=True)
tb(slide, Inches(5.2), y, Inches(7.2), Inches(0.6), "No se puede avanzar si la fase anterior no ha sido validada", size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 – RESUMEN GLOBAL DE RESULTADOS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Resultados Globales", "Accuracy por fase del torneo")

reports_data = {}
for stage_name in ["R32", "R16", "QF", "SF"]:
    try:
        with open(f"outputs/reports/reporte_{stage_name}.json", encoding="utf-8") as f:
            reports_data[stage_name] = json.load(f)
    except:
        reports_data[stage_name] = None

# Summary cards per phase
stages_info = [
    ("R32", "Dieciseisavos", 16),
    ("R16", "Octavos", 8),
    ("QF", "Cuartos", 4),
    ("SF", "Semifinales", 2),
]
x_pos = 0.5
for stage_abbr, stage_name, total in stages_info:
    data = reports_data.get(stage_abbr)
    acc = data.get("prediction_accuracy", 0) if data else 0
    hits = data.get("prediction_hits", 0) if data else 0
    validated = data.get("n_matches_validated", 0) if data else 0
    add_rounded_rect(slide, Inches(x_pos), Inches(1.5), Inches(2.95), Inches(1.5), NAVY)
    tb(slide, Inches(x_pos + 0.1), Inches(1.55), Inches(2.75), Inches(0.35), stage_abbr, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(slide, Inches(x_pos + 0.1), Inches(1.95), Inches(2.75), Inches(0.25), stage_name, size=12, color=SKY_BLUE_LIGHT, align=PP_ALIGN.CENTER)
    tb(slide, Inches(x_pos + 0.1), Inches(2.3), Inches(2.75), Inches(0.5), f"{hits}/{validated} aciertos  |  {acc:.0%}", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    x_pos += 3.2

# Evolution chart area
add_rounded_rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.8), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
tb(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(0.3), "Evolucion del accuracy por fase", size=15, color=NAVY, bold=True)
acc_data = []
stages_ordered = ["R32", "R16", "QF", "SF"]
for s in stages_ordered:
    d = reports_data.get(s)
    acc_data.append(d.get("prediction_accuracy", 0) if d else 0)

# Draw bars
bar_y = Inches(4.0)
bar_max_w = Inches(8)
bar_h = Inches(0.35)
colors = [SKY_BLUE, ACCENT_TEAL, ACCENT_GREEN, ACCENT_GOLD]
for i, (s, a) in enumerate(zip(stages_ordered, acc_data)):
    x_bar = Inches(1.0)
    y_bar_top = Inches(4.0 + i * 0.45)
    # label
    tb(slide, Inches(0.8), y_bar_top + Inches(0.03), Inches(0.8), Inches(0.3), s, size=11, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # bar bg
    add_rounded_rect(slide, x_bar, y_bar_top, bar_max_w, bar_h, SUBTLE_BG)
    # bar fill
    if a > 0:
        w = int(bar_max_w * a)
        add_rounded_rect(slide, x_bar, y_bar_top, w, bar_h, colors[i])
    # text
    tb(slide, x_bar + Inches(0.1), y_bar_top + Inches(0.02), Inches(2), Inches(0.3), f"{a:.0%}", size=11, color=WHITE if a > 0.3 else NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # details
    data_i = reports_data.get(s)
    if data_i:
        hits_i = data_i.get("prediction_hits", 0)
        total_i = data_i.get("prediction_matches", 0)
        tb(slide, x_bar + Inches(2.5), y_bar_top + Inches(0.02), Inches(4), Inches(0.3), f"{hits_i}/{total_i} partidos correctos", size=10, color=MEDIUM_TEXT, anchor=MSO_ANCHOR.MIDDLE)

# Best model per phase
add_rounded_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.6), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.03), SKY_BLUE)
tb(slide, Inches(0.8), Inches(5.55), Inches(5), Inches(0.3), "Mejor modelo por fase", size=14, color=NAVY, bold=True)
x_m = 0.8
for s in stages_ordered:
    d = reports_data.get(s)
    if d and "best_classifier" in d:
        best = d["best_classifier"]
        best_m = d.get("best_classifier_metrics", {})
        acc_best = best_m.get("accuracy", 0)
        tb(slide, Inches(x_m), Inches(5.95), Inches(2.8), Inches(0.25), f"{s}: {best} ({acc_best:.0%})", size=11, color=MEDIUM_TEXT)
        x_m += 3.1

# ═══════════════════════════════════════════════════════════
# SLIDE 10 – DETALLE R32
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Resultados - Dieciseisavos (R32)", "16 partidos  |  11 aciertos  |  68.75% accuracy")
r32 = reports_data.get("R32")
if r32:
    matches = r32.get("matches", [])
    y = Inches(1.4)
    for i, m in enumerate(matches):
        if i > 0 and i % 8 == 0:
            # continue on conceptual level - show compact
            pass
        col = i % 2
        row = i // 2
        if row > 3:
            break
        x_card = Inches(0.5 + col * 6.3)
        y_card = Inches(1.4 + row * 0.75)
        hit = m.get("hit", False)
        bg = RGBColor(0xE8, 0xF5, 0xE9) if hit else RGBColor(0xFD, 0xED, 0xED)
        add_rounded_rect(slide, x_card, y_card, Inches(6), Inches(0.65), bg, border=1, border_color=ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x_card, y_card, Inches(0.04), Inches(0.65), ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C))
        winner = m.get("predicted_winner", "?")
        actual = m.get("actual_result", {})
        real_w = actual.get("winner", "?") if actual else "?"
        sc_pred = m.get("predicted_score_90", {})
        sc_real = actual.get("score_90", {}) if actual else {}
        tb(slide, x_card + Inches(0.15), y_card + Inches(0.05), Inches(3), Inches(0.25), f"{m['team_a']} vs {m['team_b']}", size=11, color=NAVY, bold=True)
        tb(slide, x_card + Inches(0.15), y_card + Inches(0.32), Inches(4), Inches(0.25), f"Pred: {winner} ({sc_pred.get('a',0)}-{sc_pred.get('b',0)})  |  Real: {real_w} ({sc_real.get('a',0) if sc_real else '?'}-{sc_real.get('b',0) if sc_real else '?'})", size=9, color=MEDIUM_TEXT)
        # hit/miss label
        tb(slide, x_card + Inches(4.5), y_card + Inches(0.12), Inches(1.2), Inches(0.25), "ACIERTO" if hit else "FALLO", size=9, color=ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C), bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # Rest of matches compact
    add_rounded_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.4), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    tb(slide, Inches(0.8), Inches(4.7), Inches(5), Inches(0.3), "Resto de partidos R32", size=13, color=NAVY, bold=True)
    compact_text = ""
    for i, m in enumerate(matches):
        if i < 8:
            continue
        winner = m.get("predicted_winner", "?")
        actual = m.get("actual_result", {})
        real_w = actual.get("winner", "?") if actual else "?"
        hit = m.get("hit", False)
        mark = "+" if hit else "-"
        compact_text += f"{m['team_a']}-{m['team_b']}: {winner}>{real_w} ({mark})  "
    tb(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.6), compact_text, size=10, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 – DETALLE R16
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Resultados - Octavos (R16)", "8 partidos  |  6 aciertos  |  75% accuracy")
r16 = reports_data.get("R16")
if r16:
    matches = r16.get("matches", [])
    y = Inches(1.5)
    for i, m in enumerate(matches):
        col = i % 2
        row = i // 2
        x_card = Inches(0.5 + col * 6.3)
        y_card = Inches(1.5 + row * 1.35)
        hit = m.get("hit", False)
        bg = RGBColor(0xE8, 0xF5, 0xE9) if hit else RGBColor(0xFD, 0xED, 0xED)
        add_rounded_rect(slide, x_card, y_card, Inches(6), Inches(1.15), bg, border=1, border_color=ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x_card, y_card, Inches(0.04), Inches(1.15), ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C))
        tb(slide, x_card + Inches(0.15), y_card + Inches(0.08), Inches(3.5), Inches(0.3), f"{m['team_a']} vs {m['team_b']}", size=13, color=NAVY, bold=True)
        winner = m.get("predicted_winner", "?")
        prob = m.get("advance_probability_avg_a", 0)
        sc_pred = m.get("predicted_score_90", {})
        tb(slide, x_card + Inches(0.15), y_card + Inches(0.4), Inches(3.5), Inches(0.25), f"Pred: {winner} ({sc_pred.get('a',0)}-{sc_pred.get('b',0)})  Prob: {prob:.0%}", size=10, color=MEDIUM_TEXT)
        actual = m.get("actual_result", {})
        if actual:
            real_w = actual.get("winner", "?")
            sc_real = actual.get("score_90", {})
            tb(slide, x_card + Inches(0.15), y_card + Inches(0.68), Inches(3.5), Inches(0.25), f"Real: {real_w} ({sc_real.get('a',0)}-{sc_real.get('b',0)})", size=10, color=ACCENT_GREEN)
        tb(slide, x_card + Inches(4.2), y_card + Inches(0.35), Inches(1.5), Inches(0.3), "ACIERTO" if hit else "FALLO", size=11, color=ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C), bold=True, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 – DETALLE QF y SF
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Resultados - Cuartos (QF) y Semifinales (SF)", "QF: 4 partidos - 75%  |  SF: 2 partidos - pendientes")
# QF
qf = reports_data.get("QF")
if qf:
    add_rounded_rect(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(0.4), NAVY)
    tb(slide, Inches(0.8), Inches(1.42), Inches(5.5), Inches(0.35), "Cuartos de Final (QF)", size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    matches = qf.get("matches", [])
    for i, m in enumerate(matches):
        y_card = Inches(1.95 + i * 1.2)
        hit = m.get("hit", False)
        bg = RGBColor(0xE8, 0xF5, 0xE9) if hit else RGBColor(0xFD, 0xED, 0xED)
        add_rounded_rect(slide, Inches(0.5), y_card, Inches(6), Inches(1.05), bg, border=1, border_color=ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C))
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), y_card, Inches(0.04), Inches(1.05), ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C))
        tb(slide, Inches(0.8), y_card + Inches(0.08), Inches(3.5), Inches(0.25), f"{m['team_a']} vs {m['team_b']}", size=12, color=NAVY, bold=True)
        winner = m.get("predicted_winner", "?")
        sc_pred = m.get("predicted_score_90", {})
        actual = m.get("actual_result", {})
        real_w = actual.get("winner", "?") if actual else "?"
        sc_real = actual.get("score_90", {}) if actual else {}
        tb(slide, Inches(0.8), y_card + Inches(0.38), Inches(3.5), Inches(0.25), f"Pred: {winner} ({sc_pred.get('a',0)}-{sc_pred.get('b',0)})", size=10, color=MEDIUM_TEXT)
        if actual:
            tb(slide, Inches(0.8), y_card + Inches(0.65), Inches(3.5), Inches(0.25), f"Real: {real_w} ({sc_real.get('a',0)}-{sc_real.get('b',0)})", size=10, color=ACCENT_GREEN)
        tb(slide, Inches(4.5), y_card + Inches(0.35), Inches(1.5), Inches(0.3), "ACIERTO" if hit else "FALLO", size=10, color=ACCENT_GREEN if hit else RGBColor(0xE7, 0x4C, 0x3C), bold=True, align=PP_ALIGN.CENTER)

# SF
sf = reports_data.get("SF")
if sf:
    add_rounded_rect(slide, Inches(6.9), Inches(1.4), Inches(6), Inches(0.4), NAVY)
    tb(slide, Inches(7.2), Inches(1.42), Inches(5.5), Inches(0.35), "Semifinales (SF) - Pendientes", size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    matches = sf.get("matches", [])
    for i, m in enumerate(matches):
        y_card = Inches(1.95 + i * 1.2)
        add_rounded_rect(slide, Inches(6.9), y_card, Inches(6), Inches(1.05), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), y_card, Inches(0.04), Inches(1.05), SKY_BLUE)
        tb(slide, Inches(7.2), y_card + Inches(0.08), Inches(3.5), Inches(0.25), f"{m['team_a']} vs {m['team_b']}", size=12, color=NAVY, bold=True)
        winner = m.get("predicted_winner", "?")
        prob = m.get("advance_probability_avg_a", 0)
        sc_pred = m.get("predicted_score_90", {})
        tb(slide, Inches(7.2), y_card + Inches(0.38), Inches(3.5), Inches(0.25), f"Pred: {winner} ({sc_pred.get('a',0)}-{sc_pred.get('b',0)})", size=10, color=MEDIUM_TEXT)
        tb(slide, Inches(7.2), y_card + Inches(0.65), Inches(3.5), Inches(0.25), f"Prob: {prob:.0%}  |  LR: {m['advance_probability'].get('logreg',0):.0%}  RF: {m['advance_probability'].get('random_forest',0):.0%}  XGB: {m['advance_probability'].get('xgboost',0):.0%}", size=9, color=MEDIUM_TEXT)
        tb(slide, Inches(10), y_card + Inches(0.35), Inches(2.5), Inches(0.3), "PENDIENTE", size=10, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.CENTER)

# Validation metrics summary below
add_rounded_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
tb(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.3), "Accuracy acumulada del torneo:  R32 68.75%  >  R16 75%  >  QF 75%  >  SF (pendiente)", size=12, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.3), "Total: 20/28 partidos acertados en fases validadas (R32 + R16 + QF)", size=11, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 – FEATURE IMPORTANCE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Importancia de Features", "Variables mas influyentes por modelo")
top_feats = {}
for model_key in ["logreg", "random_forest", "xgboost"]:
    if r32 and model_key in r32.get("top_features", {}):
        top_feats[model_key] = r32["top_features"][model_key]

if top_feats:
    models_display = {"logreg": "Logistic Regression", "random_forest": "Random Forest", "xgboost": "XGBoost"}
    cols = [0.5, 4.8, 9.1]
    for idx, (mkey, mlabel) in enumerate(models_display.items()):
        x = cols[idx]
        add_rounded_rect(slide, Inches(x), Inches(1.5), Inches(4), Inches(0.5), NAVY)
        tb(slide, Inches(x + 0.15), Inches(1.55), Inches(3.7), Inches(0.4), mlabel, size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        feats = top_feats.get(mkey, [])
        y = Inches(2.15)
        for feat_name, imp_val in feats:
            add_rounded_rect(slide, Inches(x), y, Inches(4), Inches(0.35), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
            tb(slide, Inches(x + 0.1), y + Inches(0.03), Inches(2.2), Inches(0.28), feat_name, size=10, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
            max_val = max(abs(v) for _, v in feats) if feats else 1
            bar_w = Inches(1.2 * abs(imp_val) / max_val)
            bar_color = SKY_BLUE if imp_val > 0 else RGBColor(0xE7, 0x4C, 0x3C)
            if bar_w > Inches(0.01):
                add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(x + 2.3), y + Inches(0.08), bar_w, Inches(0.18), bar_color)
            tb(slide, Inches(x + 2.3) + bar_w + Inches(0.05), y + Inches(0.03), Inches(0.6), Inches(0.28), f"{imp_val:.3f}", size=9, color=MEDIUM_TEXT, anchor=MSO_ANCHOR.MIDDLE)
            y += Inches(0.4)
else:
    add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.3), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    tb(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.3), "Feature Importance", size=15, color=NAVY, bold=True)
    tb(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(2.5),
       "Mide la contribucion de cada variable a la decision del modelo.\n\nLR: Coeficientes de regresion (positivo > favorece a equipo A)\nRF / XGB: Importancia por ganancia de informacion\n\nFeatures mas relevantes:\nko_stage_historical_win_rate_a/b  |  key_players_prev_wc_goals_diff\nfifa_ranking_pos_diff  |  wc_best_stage_diff", size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 – DASHBOARD
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Dashboard Interactivo", "Visualizacion con Streamlit y Plotly")
sections = [
    ("Resumen", "Vista general del torneo con predicciones y accuracy por fase"),
    ("Predicciones", "Listado de partidos con marcador y probabilidad"),
    ("Modelos", "Comparacion de rendimiento entre clasificadores"),
    ("Explicabilidad", "Feature importance por fase"),
    ("Metricas", "Evolucion de accuracy, F1-score y errores"),
    ("Historial", "Predicciones historicas vs resultados oficiales"),
    ("Arbol del Mundial", "Cuadro eliminatorio interactivo"),
    ("Configuracion", "Parametros del pipeline"),
]
for i, (title, desc) in enumerate(sections):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.2)
    y = Inches(1.5 + row * 2.9)
    add_rounded_rect(slide, x, y, Inches(2.95), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(2.95), Inches(0.04), SKY_BLUE)
    tb(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.55), Inches(0.3), title, size=15, color=NAVY, bold=True)
    tb(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.55), Inches(1.5), desc, size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 – EVALUACION Y METRICAS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Evaluacion y Metricas", "Accuracy, Precision, Recall, F1-Score")
# Reports generated
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.3), "Reportes generados por fase", size=16, color=NAVY, bold=True)
reports = [
    "Reporte JSON con predicciones detalladas",
    "Validacion: accuracy, precision, recall, F1 por modelo",
    "Ranking de modelos ordenados por rendimiento",
    "Seleccion automatica del mejor clasificador",
    "Comparacion prediccion vs resultado real (hit/miss)",
    "Historial acumulado (metrics_history.csv)",
]
y = Inches(2.15)
for r in reports:
    tb(slide, Inches(0.8), y, Inches(0.2), Inches(0.25), ">", size=10, color=SKY_BLUE, bold=True)
    tb(slide, Inches(1.1), y, Inches(5.2), Inches(0.25), r, size=11, color=MEDIUM_TEXT)
    y += Inches(0.3)

# Metric definitions
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.3), "Metricas clave", size=16, color=NAVY, bold=True)
metrics_def = [
    ("Accuracy", "Porcentaje de aciertos en prediccion del ganador"),
    ("Precision", "De los que predijo como ganador, cuantos acerto"),
    ("Recall", "De los ganadores reales, cuantos predijo"),
    ("F1-Score", "Media armonica de precision y recall"),
    ("MAE / RMSE", "Error absoluto/raiz en marcadores"),
]
y = Inches(2.15)
for met, desc in metrics_def:
    tb(slide, Inches(7.2), y, Inches(1.5), Inches(0.25), met, size=11, color=NAVY, bold=True)
    tb(slide, Inches(8.8), y, Inches(3.8), Inches(0.25), desc, size=11, color=MEDIUM_TEXT)
    y += Inches(0.3)

# Validation mechanism
add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.3), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.3), "Validacion automatica", size=16, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(4.9), Inches(5.3), Inches(2.0),
   "advance_to_next_stage_guard() impide avanzar si:\n\n1. No existen resultados oficiales necesarios\n2. La fase anterior no ha sido evaluada\n\nMensaje tipico:\n[DETENIDO] La fase R16 aun no ha sido evaluada.\nEjecute: python scripts/update_stage.py R16", size=11, color=MEDIUM_TEXT)

# Weights
add_rounded_rect(slide, Inches(6.9), Inches(4.3), Inches(6), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.3), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(4.5), Inches(5.5), Inches(0.3), "Pesos dinamicos del Ensemble", size=16, color=NAVY, bold=True)
tb(slide, Inches(7.2), Inches(4.9), Inches(5.3), Inches(2.0),
   "Peso_m = accuracy_m / sum(accuracy_total)\n\nEjemplo tras validacion de R32:\nLR:  0.500  >  Peso: 0.27\nRF:  0.500  >  Peso: 0.27\nXGB: 0.833  >  Peso: 0.46\n\nEl modelo mas fiable tiene mayor influencia.", size=11, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 – ESTRUCTURA DEL PROYECTO
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Estructura del Proyecto", "Organizacion del codigo y los datos")
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
tree = (
    "mundial2026_predictor/\n"
    "  dashboard/          Aplicacion Streamlit (8 paginas)\n"
    "  data/               Datos del torneo\n"
    "    raw/              Datos originales (CSV, JSON)\n"
    "    processed/        Predicciones y comparaciones\n"
    "  outputs/            Reportes, logs e historicos\n"
    "  saved_models/       Modelos entrenados (.joblib)\n"
    "  scripts/            Scripts de ejecucion\n"
    "    run_stage.py      Ejecutar prediccion de fase\n"
    "    update_stage.py   Validar fase con resultados\n"
    "    train_models.py   Entrenar modelos\n"
    "  src/                Codigo fuente del pipeline\n"
    "    pipeline.py       Orquestador principal\n"
    "    models.py         Definicion de modelos ML\n"
    "    feature_engineering.py  Construccion de features\n"
    "    evaluation.py     Metricas y validacion\n"
    "    config.py         Configuracion central\n"
    "  tests/              Pruebas unitarias (pytest)"
)
tb(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0), tree, size=11, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 14 – COMO EJECUTAR
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Ejecucion del Pipeline", "Comandos para correr cada fase")
# Run
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.3), "Predecir una fase", size=16, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(2.1), Inches(5.3), Inches(1.6),
   "python scripts/train_models.py\n\npython scripts/run_stage.py R32\npython scripts/run_stage.py R16\npython scripts/run_stage.py QF\npython scripts/run_stage.py SF\npython scripts/run_stage.py FINAL", size=11, color=ACCENT_TEAL)
# Update
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(2.5), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.3), "Validar una fase", size=16, color=NAVY, bold=True)
tb(slide, Inches(7.2), Inches(2.1), Inches(5.3), Inches(1.6),
   "python scripts/update_stage.py R32\npython scripts/update_stage.py R16\npython scripts/update_stage.py QF\npython scripts/update_stage.py SF\npython scripts/update_stage.py FINAL\n\nstreamlit run dashboard/app.py", size=11, color=ACCENT_TEAL)
# Tests
add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.3), "Notas importantes", size=16, color=NAVY, bold=True)
tb(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(1.8),
   "Los modelos se guardan en saved_models/ con Joblib y se reutilizan entre ejecuciones.\nSi se agregan nuevos datos historicos, se pueden reentrenar con train_models.py.\n\nTests:  pytest -q  (ejecuta todas las pruebas unitarias)\n\nEl pipeline impide avanzar si la fase anterior no ha sido validada.", size=12, color=MEDIUM_TEXT)

# ═══════════════════════════════════════════════════════════
# SLIDE 15 – CONCLUSIONES
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SOFT_WHITE)
slide_header(slide, "Conclusiones y Trabajo Futuro")
# Achievements
add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(5.3), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.3), "Logros", size=17, color=NAVY, bold=True)
logros = [
    "Pipeline completo y reproducible para 6 fases",
    "Ensemble adaptable con pesos dinamicos por accuracy",
    "Validacion automatica contra resultados reales",
    "Dashboard interactivo con 8 secciones",
    "Feature engineering con 31 variables predictivas",
    "Persistencia de modelos con Joblib",
    "Cobertura de pruebas con Pytest",
]
y = Inches(2.2)
for l in logros:
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.9), y + Inches(0.02), Inches(0.08), Inches(0.08), SKY_BLUE)
    tb(slide, Inches(1.2), y, Inches(5), Inches(0.25), l, size=12, color=MEDIUM_TEXT)
    y += Inches(0.35)
# Future
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6), Inches(5.3), CARD_WHITE, border=1, border_color=SKY_BLUE_LIGHT)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(0.04), SKY_BLUE)
tb(slide, Inches(7.2), Inches(1.7), Inches(5.5), Inches(0.3), "Trabajo futuro", size=17, color=NAVY, bold=True)
futuro = [
    "Incorporar datos de lesiones en tiempo real",
    "Modelo de Deep Learning (redes neuronales)",
    "Prediccion de formacion tactica (XI inicial)",
    "Analisis de odds de casas de apuestas",
    "Simulacion Monte Carlo para probabilidades del torneo",
    "Soporte para mas torneos (Copa America, Eurocopa)",
]
y = Inches(2.2)
for f_text in futuro:
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.3), y + Inches(0.02), Inches(0.08), Inches(0.08), ACCENT_GOLD)
    tb(slide, Inches(7.6), y, Inches(5), Inches(0.25), f_text, size=12, color=MEDIUM_TEXT)
    y += Inches(0.35)

# ═══════════════════════════════════════════════════════════
# SLIDE 16 – CIERRE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.06), SKY_BLUE)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), SKY_BLUE)
# Center content
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.8), Inches(4.3), Inches(0.06), SKY_BLUE)
tb(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.8), "Gracias", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.5), "Predictor Mundial FIFA 2026", size=22, color=SKY_BLUE_LIGHT, align=PP_ALIGN.CENTER)
tb(slide, Inches(2), Inches(4.0), Inches(9), Inches(0.4), "Machine Learning aplicado al analisis deportivo", size=15, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.06), SKY_BLUE)
tb(slide, Inches(2), Inches(5.4), Inches(9), Inches(0.3), "Ciencia de Datos  |  Inteligencia Artificial  |  Deportes", size=13, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

# ─── Save ───
output_path = os.path.join(os.path.dirname(__file__), "presentacion", "Presentacion_Predictor_Mundial_2026.pptx")
prs.save(output_path)
print(f"Presentacion creada: {output_path}")
